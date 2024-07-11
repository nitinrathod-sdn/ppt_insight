import openai
from pptx import Presentation
import os
from flask import Flask, request, jsonify, render_template
from werkzeug.utils import secure_filename
import subprocess
from subprocess import run
import re
import nltk
from nltk.corpus import stopwords
from nltk.tokenize import word_tokenize
from nltk.stem import WordNetLemmatizer
from contractions import fix
from googletrans import Translator


# Set OpenAI API key
openai.api_key = os.getenv("OPENAI_API_KEY")

app = Flask(__name__)
UPLOAD_FOLDER = 'static/upload'
ALLOWED_EXTENSIONS = {'ppt', 'pptx', 'pdf', 'doc', 'docx'}

if not os.path.exists(UPLOAD_FOLDER):
    os.makedirs(UPLOAD_FOLDER)

#app.config['UPLOAD_FOLDER'] = UPLOAD_FOLDER

def allowed_file(filename):
    return '.' in filename and filename.rsplit('.', 1)[1].lower() in ALLOWED_EXTENSIONS

def extract_text_from_ppt(file_path):
    prs = Presentation(file_path)
    text_runs = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text_runs.append(shape.text)
    all_text = "\n".join(text_runs)
    return all_text

def clean_text(text):
    # Expand contractions
    text = fix(text)
    # Convert to lowercase
    text = text.lower()
    # Remove extra whitespace
    text = re.sub(r'\s+', ' ', text).strip()
    # Remove special characters and numbers
    text = re.sub(r'[^\w\s,.?!]', '', text)
    text = re.sub(r'\d+', '', text)
    # Tokenize text
    words = word_tokenize(text)
    # Remove stopwords
    stop_words = set(stopwords.words('english'))
    words = [word for word in words if word not in stop_words]
    # Lemmatize words
    lemmatizer = WordNetLemmatizer()
    words = [lemmatizer.lemmatize(word) for word in words]
    # Join words back into a single string
    cleaned_text = ' '.join(words)
    return cleaned_text 

#translation function
def translate_text(text, target_language='en'):
    translator = Translator()
    translated = translator.translate(text, dest=target_language)
    return translated.text

def generate_insights_with_gpt(text):
    response = openai.ChatCompletion.create(
        model="gpt-4o",
        messages=[
            {"role": "system", "content": "You are an assistant that provides insights and summaries of presentation content."},
            {"role": "user", "content": f"""Here is the content of a PowerPoint presentation
             always remember keep your answers more shorter and most importantly detailed
             you should generate your response in sequence that the information is being extracted it represents which field or domain
             what are the highlighted key takeaways on basis of bussiness decision is usually made:\n\n{text}\n\nPlease provide detailed insights and key takeaways from this presentation."""}
        ],
        temperature=0.7,
        max_tokens=100 # Adjust based on how much text you expect in response

    )
    return response.choices[0].message['content']


@app.route('/')
def index():
    return render_template('index.html')

@app.route('/upload', methods=['POST'])
def upload_file():
    if 'file' not in request.files:
        return jsonify(error="No file part"), 400

    file = request.files['file']
    if file.filename == '':
        return jsonify(error="No selected file"), 400

    if file and allowed_file(file.filename):
        filename = secure_filename(file.filename)
        file_path = os.path.join(app.config['UPLOAD_FOLDER'], filename)
        app.logger.info(f"File path received: {file_path}")  # Add debug logging
        
        file.save(file_path)
        
        try:
            extracted_text = extract_text_from_ppt(file_path)
            cleaned_text = clean_text(extracted_text)
            insights = generate_insights_with_gpt(cleaned_text)
        except Exception as e:
            return jsonify(error=str(e)), 500
        
        return jsonify(extracted_text=cleaned_text, insights=insights), 200
    else:
        return jsonify(error="Unsupported file format. Only .ppt and .pptx are allowed."), 400

if __name__ == "__main__":
    app.run(debug=True)
